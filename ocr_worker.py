"""
ocr_worker.py — OCR + Office document parsing sidecar (Gap 8 + Gap 9)

Handles:
  - Scanned PDFs (Tesseract OCR via pytesseract)
  - DOCX (python-docx)
  - XLSX (openpyxl)
  - PPTX (python-pptx)

All outputs match the same JSON schema as pdf_parser.py handle_parse():
  {"text": "...", "pages": N, "tables": [...], "metadata": {...}}

Install dependencies:
  pip install pytesseract pillow pdf2image python-docx openpyxl python-pptx
  # System: brew install tesseract  OR  apt-get install tesseract-ocr

Usage (from Node.js via execFile or Celery):
  python ocr_worker.py ocr      <path/to/scanned.pdf>   [language]
  python ocr_worker.py docx     <path/to/file.docx>
  python ocr_worker.py xlsx     <path/to/file.xlsx>
  python ocr_worker.py pptx     <path/to/file.pptx>
"""

import sys
import json
import os

# ---------------------------------------------------------------------------
# OCR — Scanned PDF via Tesseract (Gap 8)
# ---------------------------------------------------------------------------
def handle_ocr(file_path: str, language: str = "eng") -> dict:
    """Extract text from a scanned PDF using Tesseract OCR."""
    result = {"text": "", "pages": 0, "tables": [], "metadata": {}}

    try:
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image

        # Map BCP47 lang codes to Tesseract lang codes
        lang_map = {
            "en": "eng", "fr": "fra", "de": "deu", "es": "spa",
            "hi": "hin", "ar": "ara", "zh": "chi_sim", "ja": "jpn",
            "ru": "rus", "pt": "por", "it": "ita",
        }
        tess_lang = lang_map.get(language.split("-")[0].lower(), "eng")

        pages = convert_from_path(file_path, dpi=300)
        texts = []
        for i, page_img in enumerate(pages):
            page_text = pytesseract.image_to_string(page_img, lang=tess_lang)
            texts.append(page_text)

        result["text"] = "\n\n".join(texts)
        result["pages"] = len(pages)
        result["metadata"]["ocr_language"] = tess_lang
        result["metadata"]["ocr_engine"] = "tesseract"

    except ImportError as e:
        result["error"] = f"OCR dependencies missing: {e}. Install: pip install pytesseract pdf2image pillow"
    except Exception as e:
        result["error"] = f"OCR failed: {e}"

    return result


# ---------------------------------------------------------------------------
# DOCX Parser (Gap 9)
# ---------------------------------------------------------------------------
def handle_docx(file_path: str) -> dict:
    """Extract text and tables from a Word document."""
    result = {"text": "", "pages": 1, "tables": [], "metadata": {}}

    try:
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(file_path)
        paragraphs = []
        tables_data = []

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                text = element.text_content() if hasattr(element, "text_content") else ""
                # Use python-docx paragraph API
                pass
            elif tag == "tbl":
                pass  # handled below

        # Paragraphs (preserving headings)
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                paragraphs.append(f"\n## {para.text}\n")
            elif para.text.strip():
                paragraphs.append(para.text)

        # Tables
        for tbl in doc.tables:
            rows = []
            for row in tbl.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                tables_data.append(rows)
                # Also add table as markdown text
                header = " | ".join(rows[0])
                separator = " | ".join(["---"] * len(rows[0]))
                body = "\n".join(" | ".join(r) for r in rows[1:])
                paragraphs.append(f"\n{header}\n{separator}\n{body}\n")

        # Core properties
        props = doc.core_properties
        result["text"] = "\n\n".join(p for p in paragraphs if p.strip())
        result["tables"] = tables_data
        result["pages"] = max(1, len(paragraphs) // 30)
        result["metadata"] = {
            "title": props.title or "",
            "author": props.author or "",
            "created": str(props.created) if props.created else "",
        }

    except ImportError:
        result["error"] = "python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        result["error"] = f"DOCX parse failed: {e}"

    return result


# ---------------------------------------------------------------------------
# XLSX Parser (Gap 9)
# ---------------------------------------------------------------------------
def handle_xlsx(file_path: str) -> dict:
    """Extract text and table data from an Excel workbook."""
    result = {"text": "", "pages": 0, "tables": [], "metadata": {}}

    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_text = []
        all_tables = []

        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(cells)
            if rows:
                all_tables.append(rows)
                # Sheet as markdown table
                header = " | ".join(rows[0])
                separator = " | ".join(["---"] * len(rows[0]))
                body = "\n".join(" | ".join(r) for r in rows[1:50])  # cap at 50 rows
                all_text.append(f"## Sheet: {sheet.title}\n\n{header}\n{separator}\n{body}")

        result["text"] = "\n\n".join(all_text)
        result["tables"] = all_tables
        result["pages"] = len(wb.worksheets)
        result["metadata"] = {"sheets": [s.title for s in wb.worksheets]}

    except ImportError:
        result["error"] = "openpyxl not installed. Run: pip install openpyxl"
    except Exception as e:
        result["error"] = f"XLSX parse failed: {e}"

    return result


# ---------------------------------------------------------------------------
# PPTX Parser (Gap 9)
# ---------------------------------------------------------------------------
def handle_pptx(file_path: str) -> dict:
    """Extract text from a PowerPoint presentation."""
    result = {"text": "", "pages": 0, "tables": [], "metadata": {}}

    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {slide_num}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    tbl = shape.table
                    rows = []
                    for row in tbl.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        result["tables"].append(rows)
                        header = " | ".join(rows[0])
                        separator = " | ".join(["---"] * len(rows[0]))
                        body = "\n".join(" | ".join(r) for r in rows[1:])
                        slide_parts.append(f"\n{header}\n{separator}\n{body}")
            slides_text.append("\n".join(slide_parts))

        result["text"] = "\n\n".join(slides_text)
        result["pages"] = len(prs.slides)
        result["metadata"] = {
            "slides": len(prs.slides),
            "author": str(prs.core_properties.author or ""),
        }

    except ImportError:
        result["error"] = "python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        result["error"] = f"PPTX parse failed: {e}"

    return result


# ---------------------------------------------------------------------------
# Main router
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(json.dumps({"error": "Usage: ocr_worker.py <ocr|docx|xlsx|pptx> <file_path> [language]"}))
        sys.exit(1)

    cmd = sys.argv[1].lower()
    file_path = sys.argv[2]
    language = sys.argv[3] if len(sys.argv) > 3 else "en"

    if not os.path.exists(file_path):
        print(json.dumps({"error": f"File not found: {file_path}"}))
        sys.exit(1)

    handler_map = {
        "ocr": lambda: handle_ocr(file_path, language),
        "docx": lambda: handle_docx(file_path),
        "xlsx": lambda: handle_xlsx(file_path),
        "pptx": lambda: handle_pptx(file_path),
    }

    handler = handler_map.get(cmd)
    if not handler:
        print(json.dumps({"error": f"Unknown command: {cmd}"}))
        sys.exit(1)

    print(json.dumps(handler()))
